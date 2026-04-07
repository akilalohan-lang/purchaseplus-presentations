#!/usr/bin/env python3
"""
PurchasePlus Presentation Generator
Usage: python generate.py --config configs/ihg_chiang_mai.json
"""

import json
import argparse
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand colours
C_NAVY   = RGBColor(0x0F, 0x25, 0x42)
C_BLUE   = RGBColor(0x1A, 0x56, 0x8A)
C_TEAL   = RGBColor(0x2F, 0xAE, 0xE3)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY   = RGBColor(0xAA, 0xBB, 0xCC)
C_GREEN  = RGBColor(0x2E, 0xCC, 0x71)
C_ORANGE = RGBColor(0xE3, 0x6B, 0x2F)

# ── Slide dimensions (widescreen 16:9)
W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color=C_NAVY):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=24, color=C_WHITE, bold=False, italic=False,
             align=PP_ALIGN.LEFT, word_wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=18, color=C_WHITE, bold=False, line_spacing=Pt(6), font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return txBox


def topic_badge(slide, number, optional=False):
    """Small topic number badge top-left."""
    label = f"TOPIC {number}" + (" — IF TIME" if optional else "")
    add_rect(slide, Inches(0.4), Inches(0.3), Inches(1.6), Inches(0.35), fill_color=C_BLUE)
    add_text(slide, label, Inches(0.4), Inches(0.3), Inches(1.6), Inches(0.35),
             font_size=9, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)


def pp_logo_text(slide):
    """PurchasePlus logo text — top right."""
    add_text(slide, "PurchasePlus", Inches(10.8), Inches(0.2), Inches(2.3), Inches(0.4),
             font_size=12, color=C_TEAL, bold=True, align=PP_ALIGN.RIGHT)


def bottom_bar(slide, text=""):
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), fill_color=C_BLUE)
    if text:
        add_text(slide, text, Inches(0.4), Inches(7.1), Inches(12), Inches(0.4),
                 font_size=10, color=C_WHITE, align=PP_ALIGN.LEFT)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs, data):
    slide = blank_slide(prs)
    bg(slide, C_NAVY)

    # Left accent bar
    add_rect(slide, 0, 0, Inches(0.12), H, fill_color=C_TEAL)

    # Top blue bar
    add_rect(slide, 0, 0, W, Inches(0.08), fill_color=C_BLUE)

    # Main title
    add_text(slide, data["title"], Inches(0.5), Inches(1.8), Inches(8), Inches(1.5),
             font_size=52, bold=True, color=C_WHITE)

    # Teal underline
    add_rect(slide, Inches(0.5), Inches(3.4), Inches(4), Inches(0.06), fill_color=C_TEAL)

    # Subtitle
    add_text(slide, data["subtitle"], Inches(0.5), Inches(3.6), Inches(9), Inches(0.8),
             font_size=24, color=C_TEAL, bold=True)

    # Meta
    add_text(slide, data["meta"], Inches(0.5), Inches(4.5), Inches(9), Inches(0.5),
             font_size=14, color=C_GREY, italic=True)

    # Right side — IHG reference
    add_rect(slide, Inches(9.8), Inches(2.5), Inches(3.1), Inches(2), fill_color=C_BLUE)
    add_text(slide, "IHG OA & Finance\nConference 2026", Inches(9.9), Inches(2.7), Inches(2.9), Inches(1.6),
             font_size=18, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

    bottom_bar(slide, "Confidential — PurchasePlus × IHG")


def slide_section(prs, data):
    slide = blank_slide(prs)
    bg(slide, C_NAVY)

    topic_badge(slide, data["topic_number"])
    pp_logo_text(slide)

    # Left blue accent bar
    add_rect(slide, 0, 0, Inches(0.08), H, fill_color=C_BLUE)

    # Topic title
    add_text(slide, data["title"], Inches(0.4), Inches(1.0), Inches(12), Inches(1.0),
             font_size=38, bold=True, color=C_WHITE)

    # Teal rule
    add_rect(slide, Inches(0.4), Inches(2.1), Inches(5), Inches(0.05), fill_color=C_TEAL)

    # Body text
    add_text(slide, data["body"], Inches(0.4), Inches(2.3), Inches(12.4), Inches(4.2),
             font_size=20, color=C_GREY, word_wrap=True)

    bottom_bar(slide)


def slide_alignment(prs, data):
    slide = blank_slide(prs)
    bg(slide, C_NAVY)

    topic_badge(slide, data["topic_number"])
    pp_logo_text(slide)

    add_text(slide, data["title"], Inches(0.4), Inches(0.8), Inches(12), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.4), Inches(1.65), Inches(12.4), Inches(0.05), fill_color=C_TEAL)

    items = data["items"]
    # Header row
    header = items[0]
    add_rect(slide, Inches(0.4), Inches(1.8), Inches(6.0), Inches(0.5), fill_color=C_BLUE)
    add_rect(slide, Inches(6.5), Inches(1.8), Inches(6.4), Inches(0.5), fill_color=C_TEAL)
    add_text(slide, header[0], Inches(0.5), Inches(1.8), Inches(5.8), Inches(0.5),
             font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, header[1], Inches(6.5), Inches(1.8), Inches(6.3), Inches(0.5),
             font_size=14, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    row_colors = [RGBColor(0x1A, 0x2E, 0x4A), RGBColor(0x0F, 0x25, 0x42)]
    for i, item in enumerate(items[1:]):
        top = Inches(2.4 + i * 0.9)
        fill = row_colors[i % 2]
        add_rect(slide, Inches(0.4), top, Inches(6.0), Inches(0.75), fill_color=fill)
        add_rect(slide, Inches(6.5), top, Inches(6.4), Inches(0.75), fill_color=fill)
        add_text(slide, item[0], Inches(0.5), top, Inches(5.8), Inches(0.75),
                 font_size=16, color=C_GREY, align=PP_ALIGN.LEFT)
        # Green tick + value
        add_text(slide, f"✓  {item[1]}", Inches(6.5), top, Inches(6.3), Inches(0.75),
                 font_size=16, color=C_GREEN, bold=True, align=PP_ALIGN.LEFT)

    bottom_bar(slide)


def slide_video(prs, data):
    """Slide with key points left, video placeholder right."""
    slide = blank_slide(prs)
    bg(slide, C_NAVY)

    optional = data.get("optional", False)
    topic_badge(slide, data["topic_number"], optional)
    pp_logo_text(slide)

    # Title
    add_text(slide, data["title"], Inches(0.4), Inches(0.8), Inches(7), Inches(0.8),
             font_size=34, bold=True, color=C_WHITE)
    add_text(slide, data["subtitle"], Inches(0.4), Inches(1.6), Inches(7), Inches(0.5),
             font_size=16, color=C_TEAL, italic=True)
    add_rect(slide, Inches(0.4), Inches(2.15), Inches(3), Inches(0.05), fill_color=C_TEAL)

    # Key points
    for i, point in enumerate(data["key_points"]):
        top = Inches(2.4 + i * 0.9)
        add_rect(slide, Inches(0.4), top + Inches(0.15), Inches(0.06), Inches(0.45), fill_color=C_TEAL)
        add_text(slide, point, Inches(0.6), top, Inches(5.8), Inches(0.85),
                 font_size=15, color=C_WHITE, word_wrap=True)

    # Video placeholder box
    add_rect(slide, Inches(6.8), Inches(0.7), Inches(6.1), Inches(6.1),
             fill_color=RGBColor(0x1A, 0x2E, 0x4A),
             line_color=C_TEAL, line_width=Pt(2))

    vid_file = data.get("video_file", "")
    add_text(slide, "▶  VIDEO", Inches(7.8), Inches(2.8), Inches(4), Inches(0.6),
             font_size=28, color=C_TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, vid_file, Inches(7.3), Inches(3.5), Inches(5), Inches(0.5),
             font_size=12, color=C_GREY, align=PP_ALIGN.CENTER, italic=True)
    add_text(slide, "Insert > Video > This Device", Inches(7.3), Inches(4.2), Inches(5), Inches(0.4),
             font_size=11, color=RGBColor(0x55, 0x66, 0x77), align=PP_ALIGN.CENTER, italic=True)

    bottom_bar(slide)


def slide_closing(prs, data):
    slide = blank_slide(prs)
    bg(slide, C_NAVY)

    add_rect(slide, 0, 0, Inches(0.12), H, fill_color=C_TEAL)
    add_rect(slide, 0, 0, W, Inches(0.08), fill_color=C_BLUE)

    add_text(slide, data["title"], Inches(0.5), Inches(1.5), Inches(12), Inches(1.2),
             font_size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(4), Inches(2.8), Inches(5.3), Inches(0.06), fill_color=C_TEAL)

    add_text(slide, data["contact"], Inches(0.5), Inches(3.2), Inches(12), Inches(0.6),
             font_size=22, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, data["email"], Inches(0.5), Inches(3.9), Inches(12), Inches(0.5),
             font_size=18, color=C_GREY, align=PP_ALIGN.CENTER)
    add_text(slide, data["website"], Inches(0.5), Inches(4.5), Inches(12), Inches(0.5),
             font_size=18, color=C_TEAL, align=PP_ALIGN.CENTER)

    add_text(slide, "PurchasePlus", Inches(0.5), Inches(6.0), Inches(12), Inches(0.5),
             font_size=14, color=C_GREY, align=PP_ALIGN.CENTER)

    bottom_bar(slide, "purchaseplus.com")


# ── Main ──────────────────────────────────────────────────────────────────────

BUILDERS = {
    "title":     slide_title,
    "section":   slide_section,
    "alignment": slide_alignment,
    "video_slide": slide_video,
    "closing":   slide_closing,
}


def generate(config_path):
    with open(config_path) as f:
        config = json.load(f)

    prs = new_prs()

    for slide_data in config["slides"]:
        slide_type = slide_data["type"]
        builder = BUILDERS.get(slide_type)
        if builder:
            builder(prs, slide_data)
        else:
            print(f"  Unknown slide type: {slide_type}")

    os.makedirs("output", exist_ok=True)
    out = os.path.join("output", config["output_filename"])
    prs.save(out)
    print(f"\n✓ Saved: {out}")
    return out


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PurchasePlus Presentation Generator")
    parser.add_argument("--config", default="configs/ihg_chiang_mai.json",
                        help="Path to presentation config JSON")
    args = parser.parse_args()
    generate(args.config)
